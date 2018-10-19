"""Renders PowerPoint presentations easily with Python"""

import re
import os
import webbrowser
from functools import singledispatch
from typing import Union
from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from pptx.shapes.base import BaseShape
from pptx.shapes.table import Table
from pptx.text.text import TextFrame
from pandas import DataFrame


# TEMPLATE FUNCTIONS

def open_template(template_name: str) -> Presentation:
    """
    Opens a pptx file given the template_name and returns it.

    Parameters
    ----------
    template_name: str
        The name of the file to be open.

    Returns
    -------
    pptx.presentation.Presentation
        Presentation object

    Examples
    --------
    Open a ppt template for later use

    >>> open_template('template.pptx')
    <pptx.presentation.Presentation at ...>
    """
    return Presentation(template_name)


def render_template(template_name: str, values: dict) -> Presentation:
    """
    Returns a rendered presentation given the template name and values to be rendered.

    Parameters
    ----------
    template_name: str
        Name of the presentation to be rendered.

    values: dict
        Dictionary with the values to render on the template.

    Returns
    -------
    pptx.presentation.Presentation
        Rendered presentation

    Examples
    --------

    Render a template.

    >>> values = {'presentation_title': "My Cool Presentation"}
    >>> render_template('template.pptx', values)
    <pptx.presentation.Presentation at ...>
    """
    return render_ppt(open_template(template_name), values)


def render_and_save_template(template_name: str, values: dict, filename: str) -> None:
    """
    Renders and save a presentation with the given template name,
    values to be rendered, and filename.

    Parameters
    ----------
    template_name: str
        Name of the presentation to be saved.

    values: dict
        Dictionary with the values to render on the template.

    filename: str
        Name of the file to be saved.

    Examples
    --------

    Render and save a template

    >>> values = {'presentation_title': "My Cool Presentation"}
    >>> render_and_save_template('template.pptx', values, 'presentation.pptx')
    """
    save_ppt(render_template(template_name, values), filename)


# PRESENTATION FUNCTIONS

def _create_empty_values(shape) -> Union[str, list, dict]:
    if is_paragraph(shape):
        placeholders = _get_placeholders(shape.text_frame)

        if placeholders:
            return {keyword: '' for keyword in placeholders}
        return ''
    elif is_table(shape):
        return [[None for _ in row.cells] for row in shape.table.rows]

    elif is_chart(shape):
        return {'title': "", 'data': [], 'categories': []}
    else:
        return ''


def _is_default_name(name: str) -> bool:
    _usual_default_words = {
        'Title',
        'Placeholder',
        'Connector',
        'Elbow', 'Up',
        'Left',
        'Right',
        'Down',
        'Subtitle'
    }

    if name.istitle():
        return (set(name.split()) & _usual_default_words) == {}
    return False


def get_shapes(prs: Presentation, get_all=False) -> dict:
    """
    Returns a dictionary to be filled with the values to be rendered in the presentation.

    Parameters
    ----------
    prs: pptx.presentation.Presentation
        Presentation to get the shapes from.

    get_all: bool
        If False, filters out the shapes that are believed to be created automatically

    Returns
    -------
    dict:
        Dictionary with the shape names as keys and an empty data structure to fill the values.

    Examples
    --------

    Get all the shapes in a presentation.

    >>> prs = open_template('template.pptx')
    >>> get_shapes(prs)
    {'chart': {'categories': [], 'data': [], 'title': ''},
     'client_name': '',
     'presentation_title': '',
     'slide_text': {'cpc_change': '', 'year': ''},
     'slide_title': '',
     'table': [[None, None, None], [None, None, None], [None, None, None]]}

    >>> get_shapes(prs, all=True)
    {'Title 1': '',
     'chart': {'categories': [], 'data': [], 'title': ''},
     'client_name': '',
     'presentation_title': '',
     'slide_text': {'cpc_change': '', 'year': ''},
     'slide_title': '',
     'table': [[None, None, None], [None, None, None], [None, None, None]]}
    """
    if get_all:
        return {shape.name: _create_empty_values(shape)
                for slide in prs.slides
                for shape in slide.shapes}

    return {shape.name: _create_empty_values(shape)
            for slide in prs.slides
            for shape in slide.shapes
            if not _is_default_name(shape.name)}


def get_shapes_by_name(prs: Presentation, name: str) -> list:
    """
    Returns a list of shapes with the given name in the given presentation.
    Parameters
    ----------
    prs: pptx.presentation.Presentation
        Presentation to be saved.

    name: str
        Name of the shape(s) to be returned.

    Examples
    --------

    Get all the shapes named 'chart' in a presentation.

    >>> prs = open_template('template.pptx')
    >>> get_shapes_by_name(prs, 'chart')
    [<pptx.shapes.placeholder.PlaceholderGraphicFrame at ...>]

    """
    return [shape for slide in prs.slides for shape in slide.shapes if shape.name == name]


def render_ppt(prs: Presentation, values: dict) -> Presentation:
    """
    Returns a rendered presentation given the template name and values to be rendered.

    Parameters
    ----------
    prs: pptx.presentation.Presentation
        Presentation to be rendered.

    values: dict
        Dictionary with the values to render on the template.

    Returns
    -------
    pptx.presentation.Presentation
        Rendered presentation

    Examples
    --------

    Open a template and render it.

    >>> prs = open_template('template.pptx')
    >>> values = {'presentation_title': "My Cool Presentation"}
    >>> render_ppt(prs, values)
    <pptx.presentation.Presentation at ...>
    """

    # checks each given item
    for key, value in values.items():

        # gets all the instances of the item in the presentation
        for shape in get_shapes_by_name(prs, key):

            # depending on what kind of item it is, it renders it
            if is_table(shape):
                render_table(value, shape.table)
            elif is_paragraph(shape):
                render_paragraph(value, shape.text_frame)
            elif is_chart(shape):
                render_chart(value, shape.chart)
            #except:  # pylint: disable=bare-except
            #    print(f"Failed to render {get_shape_type(shape)} object with key {key}")

    return prs


def render_and_save_ppt(prs: Presentation, values: dict, filename: str) -> None:
    """
    Renders and save a presentation with the given template name,
    values to be rendered, and filename.

    Parameters
    ----------
    prs: Presentation
        Name of the presentation to be saved.

    values: dict
        Dictionary with the values to render on the template.

    filename: str
        Name of the file to be saved.

    Examples
    --------

    Open a template, render it and save it.

    >>> values = {'presentation_title': "My Cool Presentation"}
    >>> prs = open_template('template.pptx')
    >>> render_and_save_ppt(prs, values, 'presentation.pptx')
    """
    save_ppt(render_ppt(prs, values), filename)


def save_ppt(prs: Presentation, filename: str) -> None:
    """
    Saves the given presentation with the given filename.

    Parameters
    ----------
    prs: pptx.presentation.Presentation
        Presentation to be saved.

    filename: str
        Name of the file to be saved.

    Examples
    --------

    Render a template and save it.

    >>> values = {'presentation_title': "My Cool Presentation"}
    >>> rendered_prs = render_template('template.pptx', values)
    >>> save_ppt(rendered_prs, 'presentation.pptx')
    """
    with open(filename, 'wb') as file:
        prs.save(file)


# SHAPE FUNCTIONS


def get_shape_type(shape: BaseShape) -> str:
    """
    Returns a string with the kind of the given shape.

    Parameters
    ----------
    shape: pptx.shapes.BaseShape
        Shape to get the type from.

    Returns
    -------
    string:
        String representing the type of the shape

    Examples
    --------

    Get the type of a shape.

    >>> prs = open_template('template.pptx')
    >>> shapes = get_shapes_by_name(prs, 'client_name')
    >>> get_shape_type(shapes[0])
    'paragraph'
    """
    if is_paragraph(shape):
        return 'paragraph'
    elif is_table(shape):
        return 'table'
    elif is_chart(shape):
        return 'chart'
    return ''


def is_chart(shape: BaseShape) -> bool:
    """
    Checks whether the given shape is a chart.

    Parameters
    ----------
    shape: pptx.shapes.base.BaseShape
        Shape to get whether is a chart or not.

    Returns
    -------
    bool:
        Boolean representing whether the given shape is a table or no.

    Examples
    --------

    Check whether the given shape is a table.

    >>> prs = open_template('template.pptx')
    >>> shapes = get_shapes_by_name(prs, 'client_name')
    >>> is_chart(shapes[0])
    True
    """
    return shape.has_chart


def is_paragraph(shape: BaseShape) -> bool:
    """
    Checks whether the given shape is a paragraph.

    Parameters
    ----------
    shape: pptx.shapes.base.BaseShape
        Shape to get whether is a paragraph or not.

    Returns
    -------
    bool:
        Boolean representing whether the given shape is a table or no.

    Examples
    --------

    Check whether the given shape is a table.

    >>> prs = open_template('template.pptx')
    >>> shapes = get_shapes_by_name(prs, 'client_name')
    >>> is_paragraph(shapes[0])
    True
    """
    return shape.has_text_frame


def is_table(shape: BaseShape) -> bool:
    """
    Checks whether the given shape is a table.

    Parameters
    ----------
    shape: pptx.shapes.base.BaseShape
        Shape to get whether is a table or not.

    Returns
    -------
    bool:
        Boolean representing whether the given shape is a table or no.

    Examples
    --------

    Check whether the given shape is a table.

    >>> prs = open_template('template.pptx')
    >>> shapes = get_shapes_by_name(prs, 'client_name')
    >>> is_table(shapes[0])
    False
    """
    return shape.has_table


# RENDER FUNCTIONS

@singledispatch
def render_chart(values: Union[dict, DataFrame], chart: Chart) -> None:  # pylint: disable=unused-argument
    """
    Renders the given values into the given chart.

    Parameters
    ----------
    values: dict or pandas.DataFrame
        Values to render the chart.

    chart: pptx.chart.chart.Chart
        Chart object to be rendered.

    Examples
    --------

    Render a chart from a dictionary

    >>> prs = open_template('template.pptx')
    >>> chart_values = {
    ...        'title': "My Cool Graph",
    ...        'categories': ['d1', 'd2', 'd3'],
    ...        'data':{
    ...            'displays': [500, 750, 600],
    ...            'clicks': [150, 250, 200]
    ...        }
    ...    }
    >>> shapes = get_shapes_by_name(prs, 'chart')
    >>> shape = shapes[0]
    >>> render_chart(chart_values, shape.chart)

    Render a chart for a pandas DataFrame

    >>> prs = open_template('template.pptx')
    >>> data = [
    ...     [250, 500],
    ...     [150, 750],
    ...     [350, 600],
    ...     [300, 450],
    ...     [175, 500],
    ...     [275, 700],
    ...     [125, 550],
    ... ]
    >>> pd_chart = pd.DataFrame(data,
    ...                         index=['day1', 'day2', 'day3', 'day4', 'day5', 'day6', 'day7'],
    ...                         columns=['clicks', 'displays'])
    >>> pd_chart
      clicks  displays
    0    250       500
    1    150       750
    2    350       600
    3    300       500
    4    175       500
    5    275       700
    6    125       550
    >>> pd_chart.title = "Cool Graph"
    >>> shapes = get_shapes_by_name(prs, 'chart')
    >>> shape = shapes[0]
    >>> render_chart(pd_chart, shape.chart)
    """
    raise NotImplementedError(f"Method not implemented for {type(values)} object type")


@render_chart.register(dict)
def _(values: DataFrame, chart: Chart) -> None:
    """Renders into the given chart the values in the DataFrame."""

    chart_data = CategoryChartData()

    chart_data.categories = values['categories']

    for label, series in values['data'].items():
        chart_data.add_series(label, series)

    if 'title' in values:
        chart.chart_title.text_frame.text = values['title']

    chart.replace_data(chart_data)


@render_chart.register(DataFrame)
def _(values: DataFrame, chart: Chart) -> None:
    """Renders into the given chart the values in the DataFrame."""

    chart_data = CategoryChartData()

    chart_data.categories = list(values.index)

    for label in list(values):
        chart_data.add_series(label, list(values[label].values))

    if hasattr(values, 'title'):
        chart.chart_title.text_frame.text = values.title

    chart.replace_data(chart_data)


@singledispatch
def render_paragraph(values, text_frame: TextFrame) -> None:
    """
    In the case you want to replace the whole text.

    Parameters
    ----------
    values: dict, str, int or float
        Values to render the text.

    text_frame: pptx.text.text.TextFrame
        TextFrame object to be rendered.

    Examples
    --------

    Replace full text.

    >>> prs = open_template('template.pptx')
    >>> paragraph = {
    ...     'slide_title': "Cool insight",
    ... }
    >>> shapes = get_shapes_by_name(prs, 'slide_title')
    >>> shape = shapes[0]
    >>> render_table(paragraph, shape.text_frame)


    Replace placeholders within text.

    >>> prs = open_template('template.pptx')
    >>> paragraph = {
    ...     'slide_text': {
    ...         'year': 2018,
    ...         'cpc_change': 50
    ...     }
    ... }
    >>> shapes = get_shapes_by_name(prs, 'slide_text')
    >>> shape = shapes[0]
    >>> render_table(paragraph, shape.text_frame)
    """
    paragraph = text_frame.paragraphs[0]
    p = paragraph._p  # pylint: disable=protected-access,invalid-name
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)  # pylint: disable=protected-access
    try:
        paragraph.runs[0].text = values
    except IndexError:
        paragraph.text = values


@render_paragraph.register(dict)
def _(values: dict, text_frame: TextFrame) -> None:
    """In the case you want to replace placeholders within the paragraph"""
    for paragraph in text_frame.paragraphs:
        new_text_template = paragraph.text
        keywords = re.findall(r"\{(\w+)\}", new_text_template)
        if keywords:
            new_text = new_text_template.format(**{k: values[k] for k in keywords})
            p = paragraph._p  # pylint: disable=protected-access,invalid-name
            for idx, run in enumerate(paragraph.runs):
                if idx == 0:
                    continue
                p.remove(run._r)  # pylint: disable=protected-access
            paragraph.runs[0].text = new_text


def _get_placeholders(text_frame: TextFrame) -> list:
    placeholders = []

    for paragraph in text_frame.paragraphs:
        new_text_template = paragraph.text
        keywords = re.findall(r"\{(\w+)\}", new_text_template)
        if keywords:
            placeholders.extend(keywords)

    return placeholders


@singledispatch
def render_table(values: Union[dict, list, DataFrame], table: Table) -> None:  # pylint: disable=unused-argument
    """
    Renders a table with the given values.

    Parameters
    ----------
    values: dict, list or pandas.DataFrame
        Values to render the table

    table: pptx.shapes.table.Table
        Table object to be rendered.

    Examples
    --------

    Render table from python list

    >>> prs = open_template('template.pptx')
    >>> table: [
    ...     ['header1', 'header2', 'header3'],
    ...     ['cell1', 'cell2', 'cell3'],
    ...     ['cell4', 'cell5', 'cell6']
    ... ]
    >>> shapes = get_shapes_by_name(prs, 'table')
    >>> shape = shapes[0]
    >>> render_table(table, shape.table)


    Render table from pandas DataFrame without header

    >>> prs = open_template('template.pptx')
    >>> data = [
    ...     ['header', 'header2', 'header3'],
    ...     ['cell1', 'cell2', 'cell3'],
    ...     ['cell4', 'cell5', 'cell6']
    ... ]
    >>> table_df = pd.DataFrame(data)
    >>> table_df
        col1     col2     col3
    0   header1  header2  header3
    1   cell1    cell2    cell3
    2   cell4    cell5    cell6
    >>> shapes = get_shapes_by_name(prs, 'table')
    >>> shape = shapes[0]
    >>> render_chart(table_df, shape.table)


    Render a table from a pandas DataFrame with header.

    >>> prs = open_template('template.pptx')
    >>> data = [
    ...     ['cell1', 'cell2', 'cell3'],
    ...     ['cell4', 'cell5', 'cell6']
    ... ]
    >>> table_df = pd.DataFrame(data, columns=['header', 'header2', 'header3'])
    >>> table_df
        header1  header2  header3
    0   cell1    cell2    cell3
    1   cell4    cell5    cell6
    >>> table_df.header = True
    >>> shapes = get_shapes_by_name(prs, 'table')
    >>> shape = shapes[0]
    >>> render_chart(table_df, shape.table)
    """
    raise NotImplementedError


@render_table.register(dict)
def _(values: dict, table: Table) -> None:
    """In the case you want to render placeholders within the table,
    it will call render_paragraph for each cell"""
    for row in table.rows:
        for cell in row.cells:
            render_paragraph(values, cell.text_frame)


@render_table.register(DataFrame)
def _(values: DataFrame, table: Table) -> None:
    table_rows = iter(table.rows)

    if hasattr(values, 'header') and values.header:
        for values_cell, table_cell in zip(list(values), next(table_rows).cells):
            render_paragraph(str(values_cell), table_cell.text_frame)

    for values_row, table_row in zip(list(values.values), table_rows):
        for values_cell, table_cell in zip(list(values_row), table_row.cells):
            render_paragraph(str(values_cell), table_cell.text_frame)


@render_table.register(list)
def _(values: list, table: Table) -> None:
    """In the case you want to replace the whole table,
    it will set the value for each cell in the list"""
    for values_row, table_row in zip(values, table.rows):
        for values_cell, table_cell in zip(values_row, table_row.cells):
            render_paragraph(values_cell, table_cell.text_frame)


# DOCUMENTATION FUNCTIONS

def pypyt_doc():
    """Opens pypyt's documentation in the browser."""
    webbrowser.open('https://pypyt.readthedocs.io/en/latest/index.html')
