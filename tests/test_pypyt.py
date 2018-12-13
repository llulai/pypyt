"""tests for pypyt"""
import pandas as pd

from pytest import fixture, raises
from pptx.presentation import Presentation
from pypyt import get_shape_type, get_shapes, open_template, get_shapes_by_name, render_chart


class FakePresentation:  # pylint: disable=too-few-public-methods
    """fake class to test presentations"""
    def __init__(self, slides: list):
        self.slides = slides


class FakeSlide:  # pylint: disable=too-few-public-methods
    """fake class to test slides"""
    def __init__(self, shapes: list):
        self.shapes = shapes


class FakeShape:  # pylint: disable=too-few-public-methods
    """fake class to test shapes"""
    def __init__(self, name: str):
        self.name = name
        self.has_table = False
        self.has_text_frame = False
        self.has_chart = False


class FakeParagraph:  # pylint: disable=too-few-public-methods
    """fake class to test paragraphs"""
    def __init__(self, text: str):
        self.text = text

    def __eq__(self, other):
        return self.text == other.text


class FakeTextFrame:  # pylint: disable=too-few-public-methods
    """fake class to test text frames"""
    def __init__(self, paragraphs: list):
        self.paragraphs = paragraphs

    def __eq__(self, other):
        if len(self.paragraphs) == len(other.pargraphs):
            for par1, par2 in zip(sorted(self.paragraphs), sorted(other.pargraphs)):
                if par1 != par2:
                    return False
        else:
            return False
        return True


class FakeParagraphShape(FakeShape):  # pylint: disable=too-few-public-methods
    """fake class to test paragraph shapes"""
    def __init__(self, name: str, text_frame: FakeTextFrame):
        super().__init__(name)
        self.text_frame = text_frame
        self.has_text_frame = True

    def __eq__(self, other):
        return self.name == other.name and self.text_frame == other.text_frame


class FakeRows:  # pylint: disable=too-few-public-methods
    """fake class to test rows"""
    def __init__(self, cells: list):
        self.cells = cells

    def __eq__(self, other):
        return all(el1 == el2 for el1, el2 in zip(self.cells, other.cells))\
               and len(self.cells) == len(other.cells)


class FakeTable:  # pylint: disable=too-few-public-methods
    """fake class to test tables"""
    def __init__(self, table: list):
        self.rows = [FakeRows(row) for row in table]

    def __eq__(self, other):
        return (all(row1 == row2 for row1, row2 in zip(self.rows, other.rows))
                and len(self.rows) == len(other.rows))


class FakeTableShape(FakeShape):  # pylint: disable=too-few-public-methods
    """fake class to test table shapes"""
    def __init__(self, name, table: list):
        super().__init__(name)
        self.table = FakeTable(table)
        self.has_table = True

    def __eq__(self, other):
        return self.name == other.name and self.table == other.table


class Dummy:  # pylint: disable=too-few-public-methods
    """dummy class for tests"""
    text_frame = None
    categories = None
    series = None


class FakeChart(FakeShape):  # pylint: disable=too-few-public-methods
    """fake class to test chart shapes"""
    def __init__(self, name):
        super().__init__(name)
        self.chart_data = Dummy()
        self.chart_title = Dummy()
        self.chart_title.text_frame = Dummy()
        self.chart_title.text_frame.text = None
        self.has_chart = True

    def __eq__(self, other):
        return (self.chart_data.categories == other.chart_data.categories
                and self.chart_data.series == other.chart_data.series
                and self.chart_title.text_frame.text == other.chart_title.text_frame.text)

    def __repr__(self):
        return f"{{categories: {self.chart_data.categories}," \
               f" series: {self.chart_data.series}," \
               f" title: {self.chart_title.text_frame.text}}}"

    def replace_data(self, chart_data):
        """implemented fake method to replace data in chart"""
        self.chart_data.categories = [cat.label for cat in chart_data.categories]
        self.chart_data.series = {s.name: s.values for s in chart_data._series}  # pylint: disable=protected-access


@fixture
def fake_presentation(fake_paragraph,  # pylint: disable=redefined-outer-name
                      fake_title,  # pylint: disable=redefined-outer-name
                      fake_paragraph_placeholder,  # pylint: disable=redefined-outer-name
                      fake_table,  # pylint: disable=redefined-outer-name
                      fake_chart):  # pylint: disable=redefined-outer-name
    """fake presentation instance for tests"""
    return FakePresentation([
        FakeSlide([fake_paragraph_placeholder, fake_table]),
        FakeSlide([fake_table, fake_chart]),
        FakeSlide([fake_paragraph, fake_title]),
    ])


@fixture
def fake_title():
    """fake title instance for tests"""
    return FakeShape("Title 1")


@fixture
def fake_paragraph_placeholder():
    """fake paragraph instance for tests"""
    return FakeParagraphShape('shape_paragraph_placeholder',
                              FakeTextFrame([FakeParagraph('one {place}')]))


@fixture
def fake_paragraph():
    """fake paragraph without placeholder instance for tests"""
    return FakeParagraphShape('shape_paragraph', FakeTextFrame([FakeParagraph('text')]))


@fixture
def fake_table():
    """fake table instance for tests"""
    return FakeTableShape('shape_table', [[1, 2, 3], [4, 5, 6]])


@fixture
def fake_chart():
    """fake chart instance for tests"""
    return FakeChart('shape_chart')


@fixture
def chart_values_dict():
    """chart values as dict for testing rendering"""
    return {
        'categories': ['d1', 'd2', 'd3', 'd4'],
        'data': {
            'clicks': [125, 300, 250, 200],
            'displays': [500, 450, 600, 400],
        }
    }


@fixture
def chart_values_df():
    """chart values as df for testing rendering"""
    data = [
        {'clicks': 125, 'displays': 500},
        {'clicks': 300, 'displays': 450},
        {'clicks': 250, 'displays': 600},
        {'clicks': 200, 'displays': 400},
    ]

    return pd.DataFrame(data, index=['d1', 'd2', 'd3', 'd4'])


def test_open():
    """test open a presentaion"""
    assert isinstance(open_template('tests/__template__.pptx'), Presentation)


def test_get_shapes(fake_presentation):  # pylint: disable=redefined-outer-name
    """test get shapes from presentation"""

    empty_values = {
        'shape_paragraph_placeholder': {'place': ''},
        'shape_paragraph': '',
        'shape_table': [[None, None, None], [None, None, None]],
        'shape_chart': {'title': '', 'data': [], 'categories': []}
    }

    assert get_shapes(fake_presentation) == empty_values

    empty_values['Title 1'] = ''

    assert get_shapes(fake_presentation, get_all=True) == empty_values


def test_get_shape_type(fake_paragraph_placeholder, fake_chart, fake_table, fake_title):  # pylint: disable=redefined-outer-name
    """test get shapes name"""

    assert get_shape_type(fake_paragraph_placeholder) == 'paragraph'
    assert get_shape_type(fake_table) == 'table'
    assert get_shape_type(fake_chart) == 'chart'
    assert get_shape_type(fake_title) == ''


def test_get_one_shape(fake_paragraph,  # pylint: disable=redefined-outer-name
                       fake_paragraph_placeholder,  # pylint: disable=redefined-outer-name
                       fake_table,  # pylint: disable=redefined-outer-name
                       fake_chart,  # pylint: disable=redefined-outer-name
                       fake_presentation):  # pylint: disable=redefined-outer-name
    """test get one shape from the presentatino"""

    assert get_shapes_by_name(fake_presentation,
                              'shape_paragraph_placeholder') == [fake_paragraph_placeholder]
    assert get_shapes_by_name(fake_presentation, 'shape_paragraph') == [fake_paragraph]
    assert get_shapes_by_name(fake_presentation, 'shape_table') == [fake_table, fake_table]
    assert get_shapes_by_name(fake_presentation, 'shape_chart') == [fake_chart]


def test_get_no_shape(fake_presentation):  # pylint: disable=redefined-outer-name
    """test get no shape"""
    assert get_shapes_by_name(fake_presentation, 'sh4') == []


def test_get_multiple_shapes(fake_table, fake_presentation):  # pylint: disable=redefined-outer-name
    """gtest get multiple instances of the same shape"""
    assert get_shapes_by_name(fake_presentation, 'shape_table') == [fake_table, fake_table]


def test_chart_same_output_no_title(chart_values_dict, chart_values_df):  # pylint: disable=redefined-outer-name
    """test get same output for dict and df data for charts"""
    fake_chart_dict = FakeChart('fake_chart_dict')
    fake_chart_df = FakeChart('fake_chart_df')

    render_chart(chart_values_df, fake_chart_dict)
    render_chart(chart_values_dict, fake_chart_df)

    assert fake_chart_dict == fake_chart_df


def test_chart_same_output_title(chart_values_dict, chart_values_df):  # pylint: disable=redefined-outer-name
    """test get same output with title for charts of dict or df"""
    chart_values_dict['title'] = 'titulo'
    chart_values_df.title = 'titulo'

    fake_chart_dict = FakeChart('fake_chart_dict')
    fake_chart_df = FakeChart('fake_chart_df')

    render_chart(chart_values_df, fake_chart_dict)
    render_chart(chart_values_dict, fake_chart_df)

    assert fake_chart_dict == fake_chart_df


def test_chart_invalid_values():
    """test raise error with invalid values"""
    chart_invalid_data = [1, 2, 3]

    with raises(NotImplementedError):
        render_chart(chart_invalid_data, FakeChart('shape_chart'))
